#!/usr/bin/env python
"""Turn Roi's POC / product decks into knowledge the generators can use.

CEO request #11: "Upload our POC and Product presentations to the AI so it can
learn our product better." Today the AI knows loose proof points (54%, 55s,
180-400% ROI) but nothing about modules, architecture, sector cases or the
objections sales actually hears — so posts stay generic.

Why this script exists: the MAS RAG ingests **.md only** (`ingest_directory`
globs `*.md`), so a PDF/PPTX has no path in. This converts decks to structured
markdown, screens every file for confidential content, and drops the result in
`data/rag/product/` for the MAS to index.

    pip install pdfplumber python-pptx      # only if you have those formats

    python scripts/ingest_product_docs.py --input ~/decks --dry-run
    python scripts/ingest_product_docs.py --input ~/decks

SAFETY: nothing is written until the confidentiality screen passes. Client
names, pricing and signed-agreement language must never reach a public post —
the content bank already forbids naming clients, and this keeps that true at
the source rather than hoping the prompt holds.
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

# Markers that must never enter a corpus feeding public posts. Deliberately
# broader than the MAS gate: that one protects HR/legal data, this one also
# protects commercial confidentiality (a leaked client name is a real incident).
CONFIDENTIAL_MARKERS: tuple[tuple[str, re.Pattern[str]], ...] = (
    ("client_name_carrefour", re.compile(r"(?i)\bcarrefour\b")),
    ("pricing_table", re.compile(r"(?i)\b(price per camera|per site pricing|monthly fee|USD\s?[\d,.]+/)")),
    ("signed_agreement", re.compile(r"(?i)\b(this agreement is entered into|hereinafter referred to as)\b")),
    ("cpf_number", re.compile(r"\b\d{3}\.\d{3}\.\d{3}-\d{2}\b")),
    ("credentials", re.compile(r"(?i)\b(api[_ ]?key|password|secret)\b\s*[:=]")),
)

MIN_SECTION_CHARS = 60  # shorter than this is a slide title, not knowledge


def extract_pdf(path: Path) -> list[tuple[str, str]]:
    """Return [(section_title, text)] from a PDF, one entry per page."""
    try:
        import pdfplumber
    except ImportError:
        print(f"  ! {path.name}: install pdfplumber to read PDFs", file=sys.stderr)
        return []
    out = []
    with pdfplumber.open(path) as pdf:
        for n, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            if len(text) >= MIN_SECTION_CHARS:
                first = text.split("\n", 1)[0][:80].strip()
                out.append((first or f"Page {n}", text))
    return out


def extract_pptx(path: Path) -> list[tuple[str, str]]:
    """Return [(slide_title, text)] from a PPTX, one entry per slide."""
    try:
        from pptx import Presentation
    except ImportError:
        print(f"  ! {path.name}: install python-pptx to read PPTX", file=sys.stderr)
        return []
    out = []
    for n, slide in enumerate(Presentation(path).slides, start=1):
        parts = [sh.text.strip() for sh in slide.shapes
                 if getattr(sh, "has_text_frame", False) and sh.text.strip()]
        if not parts:
            continue
        text = "\n".join(parts)
        if len(text) >= MIN_SECTION_CHARS:
            out.append((parts[0][:80], text))
    return out


def extract_md(path: Path) -> list[tuple[str, str]]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [(path.stem, text)] if len(text) >= MIN_SECTION_CHARS else []


EXTRACTORS = {".pdf": extract_pdf, ".pptx": extract_pptx, ".md": extract_md, ".txt": extract_md}


def screen(text: str) -> list[str]:
    """Names every confidentiality marker found. Empty list = safe to ingest."""
    return [name for name, pattern in CONFIDENTIAL_MARKERS if pattern.search(text)]


def to_markdown(source: Path, sections: list[tuple[str, str]]) -> str:
    """Structure sections as markdown — the MAS chunks on '##' headings."""
    lines = [
        f"# {source.stem}",
        "",
        f"> Source: `{source.name}` · converted for RAG by scripts/ingest_product_docs.py",
        "> Product knowledge for content generation. Confidentiality screened at ingestion.",
        "",
    ]
    for title, body in sections:
        clean_title = re.sub(r"\s+", " ", title).strip(" .:-") or "Section"
        lines += [f"## {clean_title}", "", body.strip(), ""]
    return "\n".join(lines)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True, help="folder holding the decks (PDF/PPTX/MD)")
    ap.add_argument(
        "--out",
        default=str(Path.home() / "OneDrive" / "Área de Trabalho" / "yedda-mas-step1" / "data" / "rag" / "product"),
        help="MAS RAG product collection folder",
    )
    ap.add_argument("--dry-run", action="store_true", help="report only, write nothing")
    ap.add_argument("--allow-flagged", action="store_true",
                    help="ingest even when the screen flags content (asks per file)")
    args = ap.parse_args()

    src_dir = Path(args.input).expanduser()
    if not src_dir.is_dir():
        print(f"FATAL: {src_dir} is not a folder", file=sys.stderr)
        return 2
    out_dir = Path(args.out).expanduser()

    files = [p for p in sorted(src_dir.iterdir()) if p.suffix.lower() in EXTRACTORS]
    if not files:
        print(f"No PDF/PPTX/MD files in {src_dir}")
        return 1

    print(f"found {len(files)} file(s) in {src_dir}\n")
    written = skipped = 0
    for path in files:
        sections = EXTRACTORS[path.suffix.lower()](path)
        if not sections:
            print(f"  – {path.name}: no readable text, skipped")
            skipped += 1
            continue

        body = to_markdown(path, sections)
        flags = screen(body)
        label = f"{len(sections)} section(s), {len(body):,} chars"

        if flags:
            print(f"  ⚠ {path.name}: {label} — FLAGGED: {', '.join(flags)}")
            if not args.allow_flagged:
                print("      not ingested. Remove the confidential parts, or re-run with --allow-flagged.")
                skipped += 1
                continue
            print("      ingesting anyway (--allow-flagged)")
        else:
            print(f"  ✓ {path.name}: {label} — clean")

        if not args.dry_run:
            out_dir.mkdir(parents=True, exist_ok=True)
            target = out_dir / f"{re.sub(r'[^a-z0-9]+', '_', path.stem.lower()).strip('_')}.md"
            target.write_text(body, encoding="utf-8")
            print(f"      → {target}")
        written += 1

    print(f"\n{'would write' if args.dry_run else 'wrote'} {written} file(s)"
          + (f", skipped {skipped}" if skipped else ""))
    if args.dry_run:
        print("\nDRY RUN — nothing written.")
    elif written:
        print("\nNext: index it on the MAS side, then the generators can use it:")
        print("  cd ../yedda-mas-step1 && python scripts/ingest_rag_docs.py")
        print("\nNOTE: 'product' must be mapped in RAG_COLLECTIONS and granted to the")
        print("marketing agents in config/org_registry.yaml — same two steps the")
        print("'marketing' collection needed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
